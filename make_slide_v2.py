"""
MTC Chain 발표 슬라이드 v2
  Slide 1 : 전체 시스템 플로우 (Actor + Business Flow)
  Slide 2 : 기술 스택 (5-layer 아키텍처)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── 색상 ───────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0f, 0x1e)
PANEL   = RGBColor(0x10, 0x18, 0x30)
PANEL2  = RGBColor(0x14, 0x1e, 0x38)
BORDER  = RGBColor(0x1e, 0x3a, 0x5f)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LIGHT   = RGBColor(0xc8, 0xd4, 0xe8)
DIM     = RGBColor(0x7a, 0x8a, 0xa8)
ACCENT  = RGBColor(0x38, 0xbd, 0xf8)   # sky-400

NAVY    = RGBColor(0x0c, 0x23, 0x40)
BLUE    = RGBColor(0x1d, 0x4e, 0xd8)
GREEN   = RGBColor(0x05, 0x96, 0x69)
AMBER   = RGBColor(0xd9, 0x77, 0x06)
GRAY    = RGBColor(0x52, 0x62, 0x7a)

# layer colors (slide 2)
L_USER  = RGBColor(0x7c, 0x3a, 0xed)   # violet
L_FE    = RGBColor(0x06, 0x82, 0xf5)   # blue
L_INFRA = RGBColor(0xf4, 0x81, 0x20)   # orange
L_BC    = RGBColor(0x62, 0x7e, 0xea)   # indigo
L_STORE = RGBColor(0x11, 0x9e, 0x7e)   # teal


# ── 기본 도형 헬퍼 ──────────────────────────────────────────────────────────────
def bg_fill(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, border_color=None, border_pt=0.75):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s


def tb(slide, text, l, t, w, h,
       size=9, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def multiline_tb(slide, lines, l, t, w, h, size=9, color=WHITE, align=PP_ALIGN.LEFT):
    """lines: list of (text, bold, size_override_or_None)"""
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, bold, sz in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz if sz else size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def vbar(slide, l, t, h, color, w=0.07):
    rect(slide, l, t, w, h, fill=color)


def arrow_down(slide, cx, y1, y2, color=ACCENT, pt=1.5):
    conn = slide.shapes.add_connector(
        1, Inches(cx), Inches(y1), Inches(cx), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(pt)


def arrow_right(slide, x1, cy, x2, color=ACCENT, pt=1.5):
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(cy), Inches(x2), Inches(cy))
    conn.line.color.rgb = color
    conn.line.width = Pt(pt)


def chip(slide, text, l, t, w, h, bg, fg=WHITE, size=8):
    rect(slide, l, t, w, h, fill=bg)
    tb(slide, text, l+0.02, t+0.01, w-0.04, h-0.02,
       size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  프레젠테이션 초기화
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 : 시스템 플로우
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg_fill(s1, BG)

# ── 타이틀 바 ────────────────────────────────────────────────────────────────
rect(s1, 0, 0, 13.33, 0.68, fill=PANEL, border_color=BORDER)
tb(s1, "MTC Chain  —  전체 시스템 플로우",
   0.22, 0.05, 9, 0.38, size=19, bold=True, color=WHITE)
tb(s1, "강재 이력 관리 시스템  ·  Ethereum Sepolia  ·  OpenZeppelin AccessControl  ·  IPFS / Pinata",
   0.22, 0.40, 10, 0.24, size=8.5, color=ACCENT)
chip(s1, "Slide 1 / 2", 12.2, 0.20, 0.95, 0.26,
     RGBColor(0x1e, 0x2e, 0x50), ACCENT, 8)

# ── 섹션 헤더 ────────────────────────────────────────────────────────────────
tb(s1, "참여자", 0.14, 0.76, 2.4, 0.22,
   size=8.5, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(s1, "비즈니스 플로우 (온체인 트랜잭션 흐름)",
   2.72, 0.76, 10.4, 0.22,
   size=8.5, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# ── 세로 구분선 ───────────────────────────────────────────────────────────────
rect(s1, 2.60, 0.72, 0.02, 6.62, fill=BORDER)

# ── LEFT: 역할 카드 ───────────────────────────────────────────────────────────
ROLES = [
    ("Admin",              "DEFAULT_ADMIN_ROLE",    NAVY,  1.05),
    ("Mill\n(제강사)",      "MILL_ROLE",             BLUE,  2.05),
    ("Fabricator\n(가공사)","FABRICATOR_ROLE",       GREEN, 3.40),
    ("Integrator\n(통합사)","INTEGRATOR_ROLE",       AMBER, 5.12),
    ("Auditor\n(감사)",     "누구나 조회 가능",        GRAY,  6.18),
]

for name, role_const, col, ry in ROLES:
    rect(s1, 0.14, ry, 2.36, 0.85, fill=PANEL2, border_color=col, border_pt=1.5)
    vbar(s1, 0.14, ry, 0.85, col)
    tb(s1, name, 0.30, ry+0.06, 2.1, 0.35,
       size=10.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    tb(s1, role_const, 0.30, ry+0.46, 2.1, 0.28,
       size=7.5, color=LIGHT, italic=True, align=PP_ALIGN.LEFT)

# ── RIGHT: 플로우 박스 ────────────────────────────────────────────────────────
#  (label, fn_call, desc, y, color, height)
FLOWS = [
    (
        "① 역할 부여",
        "grantMill() / grantFabricator() / grantIntegrator()",
        "Admin이 각 지갑 주소에 역할 할당  →  AccessControl에 역할 등록  →  RoleGranted 이벤트 발생",
        1.05, NAVY, 0.83,
    ),
    (
        "② MTC 발행",
        "issueMtc(steelId, weight, ipfsCid, pdfHash)",
        "PDF 생성  →  SHA-256 해시 계산  →  Pinata(IPFS) 업로드 → CID 획득\n→ 컨트랙트 호출  →  SteelMinted 이벤트 (steelId·mill·weight·ipfsCid·pdfHash·timestamp)",
        2.05, BLUE, 0.98,
    ),
    (
        "③ 소유권 이전",
        "transferOwnership(steelId, to)",
        "ACTIVE 강재의 owner를 Fabricator 또는 Integrator로 변경\n→ SteelOwnershipTransferred 이벤트 (from·to)",
        3.18, RGBColor(0x25, 0x6a, 0xe0), 0.85,
    ),
    (
        "④-a 강재 분할",
        "splitSteel(parentId, childWeights[])",
        "부모 → SPLIT  /  자식 N개 ACTIVE 생성  (손실 허용 ≤ 10%)\n→ SteelSplit 이벤트 (parentId·childIds[]·weights)",
        4.18, GREEN, 0.85,
    ),
    (
        "④-b 강재 조합",
        "combineSteel(parentIds[], childId, weight, cid, hash)",
        "부모들 → COMBINED  /  자식 1개 ACTIVE 생성  (손실 허용 ≤ 15%)\n→ SteelCombined 이벤트 (parentIds[]·childId·weights)",
        4.18, RGBColor(0x02, 0xb0, 0x7a), 0.85,
    ),
    (
        "⑤ 사용 등록",
        "markAsUsed(steelId, productId)",
        "상태 ACTIVE → USED  (불가역)  /  productMap[productId] = steelId 저장\n→ SteelUsed 이벤트 (steelId·productId·integrator)",
        5.18, AMBER, 0.85,
    ),
    (
        "⑥ 이력 조회",
        "getSteel() / getParents() / getChildren() / getSteelByProduct()",
        "재귀 트리 탐색으로 Ancestry 시각화  /  IPFS CID로 PDF 원본 다운로드 후 pdfHash(SHA-256) 검증",
        6.18, GRAY, 0.70,
    ),
]

# ④-a 와 ④-b 는 분기 — 좌/우 나란히 배치
SPLIT_COMBINE_Y = 4.18
split_box_w     = 4.90

for item in FLOWS:
    label, fn, desc, fy, col, fh = item

    # ④ 분기 처리
    if label == "④-a 강재 분할":
        lx = 2.72
        bw = split_box_w
    elif label == "④-b 강재 조합":
        lx = 2.72 + split_box_w + 0.14
        bw = split_box_w
    else:
        lx = 2.72
        bw = 10.40

    rect(s1, lx, fy, bw, fh, fill=PANEL2, border_color=col, border_pt=1.2)
    vbar(s1, lx, fy, fh, col)

    tb(s1, label, lx+0.14, fy+0.05, bw-0.20, 0.26,
       size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    tb(s1, fn, lx+0.14, fy+0.30, bw-0.20, 0.22,
       size=7.5, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    tb(s1, desc, lx+0.14, fy+0.50, bw-0.20, fh-0.52,
       size=7.5, color=LIGHT, align=PP_ALIGN.LEFT)

# 화살표: ①→②→③→④→⑤→⑥
ARROW_YS = [1.90, 3.05, 4.05]  # ①→②, ②→③, ③→④
for ay in ARROW_YS:
    arrow_down(s1, 7.92, ay, ay+0.13)

# ③→④ 분기 화살표 (가운데서 좌/우로)
arrow_down(s1, 7.92, 4.05, 4.18)

# ④→⑤ 화살표 (두 박스 아래 중간에서)
arrow_down(s1, 7.92, 5.03, 5.18)

# ⑤→⑥
arrow_down(s1, 7.92, 6.03, 6.18)

# OR 표시 (분기)
tb(s1, "OR", 7.79, 4.15, 0.28, 0.20,
   size=8, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ── 하단 상태 머신 요약 바 ─────────────────────────────────────────────────────
rect(s1, 0, 7.08, 13.33, 0.42, fill=PANEL, border_color=BORDER)
tb(s1,
   "상태 머신:  ACTIVE  →  SPLIT  |  COMBINED  |  USED  (모든 전이 불가역)  "
   "  ·  역할 검증 실패 시  AccessControlUnauthorizedAccount  revert  "
   "  ·  온체인 pdfHash(bytes32) = SHA-256 불변 보장",
   0.18, 7.12, 13.0, 0.34,
   size=7.5, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 : 기술 스택 (5-layer 아키텍처)
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg_fill(s2, BG)

# ── 타이틀 바 ────────────────────────────────────────────────────────────────
rect(s2, 0, 0, 13.33, 0.68, fill=PANEL, border_color=BORDER)
tb(s2, "MTC Chain  —  기술 스택 아키텍처",
   0.22, 0.05, 9, 0.38, size=19, bold=True, color=WHITE)
tb(s2, "5-Layer  풀스택 블록체인 dApp  ·  Ethereum Sepolia  ·  Cloudflare Edge  ·  IPFS",
   0.22, 0.40, 10, 0.24, size=8.5, color=ACCENT)
chip(s2, "Slide 2 / 2", 12.2, 0.20, 0.95, 0.26,
     RGBColor(0x1e, 0x2e, 0x50), ACCENT, 8)

# ── 레이어 정의 ────────────────────────────────────────────────────────────────
#  각 레이어: (번호, 이름, 역할설명, 컴포넌트[], color, y, h)
LAYERS = [
    (
        "L1", "사용자 계층",
        "User Layer",
        L_USER, 0.80, 0.95,
        [
            ("MetaMask Wallet",
             "트랜잭션 서명 (EIP-155)\nSepolia ETH 가스 지불\n계정 주소 = 역할 식별자",
             0.20),
        ],
        "MetaMask",
    ),
    (
        "L2", "프론트엔드 계층",
        "Frontend Layer",
        L_FE, 1.88, 1.42,
        [
            ("React 18 + Vite 5",
             "SPA 라우팅 (React Router v6)\nTailwind CSS 3  스타일링\nVite HMR 개발 서버",
             0.20),
            ("ethers.js  v6",
             "Contract 인터페이스 연결\nABI 인코딩 / 이벤트 파싱\nprovider · signer 관리",
             0.20),
            ("React Flow  v11",
             "Ancestry Tree 시각화\n강재 이력 DAG 렌더링\n노드/엣지 인터랙티브",
             0.20),
            ("Cloudflare Pages",
             "정적 빌드 배포 (CDN)\nmtc-chain.pages.dev\nCI/CD 자동 배포",
             0.20),
        ],
        "Frontend",
    ),
    (
        "L3", "미들웨어 계층",
        "Middleware Layer  (Edge)",
        L_INFRA, 3.43, 1.00,
        [
            ("Cloudflare Workers",
             "API Proxy  (CORS 처리)\nPinata API Key 보호\n비밀키 환경변수 격리",
             0.20),
            ("Cloudflare KV",
             "productId 메타데이터 저장\nKey-Value 엣지 캐시\n전세계 분산 저장",
             0.20),
        ],
        "Edge",
    ),
    (
        "L4", "블록체인 계층",
        "Blockchain Layer",
        L_BC, 4.56, 1.60,
        [
            ("Ethereum Sepolia\n(테스트넷)",
             "EVM 호환 퍼블릭 체인\nPoS 합의 알고리즘\n블록 탐색기: Etherscan",
             0.20),
            ("MtcRegistry.sol",
             "Solidity ^0.8.24\nOpenZeppelin AccessControl v5\n커스텀 에러·이벤트 5종",
             0.20),
            ("스마트 컨트랙트\n핵심 함수",
             "issueMtc  ·  transferOwnership\nsplitSteel  ·  combineSteel\nmarkAsUsed  ·  getSteel",
             0.20),
            ("온체인 데이터 구조",
             "struct Steel  { steelId·weight\n  ipfsCid·pdfHash·status\n  mill·owner·parentIds·childIds }",
             0.20),
        ],
        "Blockchain",
    ),
    (
        "L5", "저장소 계층",
        "Storage Layer",
        L_STORE, 6.29, 0.90,
        [
            ("IPFS  /  Pinata",
             "MTC PDF 원본 분산 저장\nCID(Content Identifier) 반환\n불변성 보장 (내용 기반 주소)",
             0.20),
            ("온체인 해시 검증",
             "pdfHash(bytes32) = SHA-256\n원본 무결성 검증 가능\n위변조 불가 증명",
             0.20),
        ],
        "Storage",
    ),
]

CARD_W     = 2.82   # 카드 너비
CARD_GAP   = 0.18   # 카드 간격
LBL_W      = 1.52   # 레이어 라벨 너비
LBL_L      = 0.14
CARD_START = LBL_L + LBL_W + 0.20

for layer_id, name_ko, name_en, col, ly, lh, cards, tag in LAYERS:

    # 레이어 전체 배경
    rect(s2, 0.14, ly, 13.05, lh, fill=PANEL2, border_color=col, border_pt=1.2)

    # 레이어 라벨 배경
    rect(s2, LBL_L, ly, LBL_W, lh, fill=col)
    tb(s2, layer_id, LBL_L+0.06, ly+0.08, LBL_W-0.12, 0.28,
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s2, name_ko, LBL_L+0.06, ly+0.35, LBL_W-0.12, 0.28,
       size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s2, name_en, LBL_L+0.06, ly+0.58, LBL_W-0.12, 0.26,
       size=7.5, color=RGBColor(0xd0,0xe8,0xff), align=PP_ALIGN.CENTER, italic=True)

    # 카드 배치
    n = len(cards)
    total_w = n * CARD_W + (n-1) * CARD_GAP
    cx = CARD_START

    for card_title, card_body, _ in cards:
        card_h = lh - 0.14
        card_top = ly + 0.07

        rect(s2, cx, card_top, CARD_W, card_h,
             fill=RGBColor(0x0d, 0x14, 0x26), border_color=col, border_pt=0.6)
        vbar(s2, cx, card_top, card_h, col, w=0.05)

        tb(s2, card_title,
           cx+0.12, card_top+0.05, CARD_W-0.18, 0.26,
           size=8.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        tb(s2, card_body,
           cx+0.12, card_top+0.32, CARD_W-0.18, card_h-0.36,
           size=7.5, color=LIGHT, align=PP_ALIGN.LEFT)

        cx += CARD_W + CARD_GAP

# ── 레이어 연결 화살표 ────────────────────────────────────────────────────────
CONN_X = 0.875  # 라벨 중앙
layer_bottoms = [
    0.80 + 0.95,   # L1 bottom
    1.88 + 1.42,   # L2 bottom
    3.43 + 1.00,   # L3 bottom
    4.56 + 1.60,   # L4 bottom
]
gap_midpoints = [
    (0.80+0.95 + 1.88) / 2,
    (1.88+1.42 + 3.43) / 2,
    (3.43+1.00 + 4.56) / 2,
    (4.56+1.60 + 6.29) / 2,
]

for i, (y1, ym) in enumerate(zip(layer_bottoms, gap_midpoints)):
    # 화살표 연결선
    arrow_down(s2, CONN_X, y1, y1 + (gap_midpoints[i] - y1),
               color=RGBColor(0x3a, 0x5a, 0x90), pt=1.2)
    arrow_down(s2, CONN_X, gap_midpoints[i], gap_midpoints[i] + (layer_bottoms[i] - y1) * 0.4,
               color=RGBColor(0x3a, 0x5a, 0x90), pt=1.2)

# ── 오른쪽 범례: 데이터 흐름 요약 ────────────────────────────────────────────
legend_l = CARD_START + 4 * CARD_W + 4 * CARD_GAP + 0.08
legend_w = 13.19 - legend_l

if legend_w > 0.5:
    rect(s2, legend_l, 0.80, legend_w, 6.39, fill=PANEL, border_color=BORDER)
    tb(s2, "데이터 흐름", legend_l+0.06, 0.85, legend_w-0.12, 0.25,
       size=8, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    flow_summary = [
        "① 사용자 서명",
        "   MetaMask",
        "",
        "② 컨트랙트 호출",
        "   ethers.js",
        "",
        "③ Proxy 경유",
        "   CF Workers",
        "",
        "④ 온체인 기록",
        "   Sepolia",
        "",
        "⑤ PDF 저장",
        "   IPFS CID",
    ]
    ty = 1.16
    for line in flow_summary:
        color = ACCENT if line.startswith("①②③④⑤") or line[:1] in "①②③④⑤" else LIGHT
        tb(s2, line, legend_l+0.06, ty, legend_w-0.12, 0.20,
           size=7.5, color=color, align=PP_ALIGN.LEFT)
        ty += 0.20

# ── 하단 요약 바 ─────────────────────────────────────────────────────────────
rect(s2, 0, 7.08, 13.33, 0.42, fill=PANEL, border_color=BORDER)
tb(s2,
   "Hardhat  컴파일 · 배포  (npx hardhat deploy --network sepolia)  "
   "  ·  OpenZeppelin v5 AccessControl  상속  "
   "  ·  ethers.js v6  Contract.connect(signer)  "
   "  ·  Pinata SDK  pinFileToIPFS()  →  IPFS CID",
   0.18, 7.12, 13.0, 0.34,
   size=7.5, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  저장
# ══════════════════════════════════════════════════════════════════════════════
OUT = r"C:\Users\wltn4\Desktop\MTC_Project\MTC_Presentation_Slides.pptx"
prs.save(OUT)
print(f"저장 완료 → {OUT}")
print("슬라이드 수:", len(prs.slides))
