"""
MTC Chain — 전체 시스템 플로우 슬라이드 생성기
16:9 와이드스크린, 블록체인 수업 발표용
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0a, 0x0f, 0x1e)   # 거의 검정 (배경)
NAVY     = RGBColor(0x0c, 0x23, 0x40)   # Admin
BLUE     = RGBColor(0x1d, 0x4e, 0xd8)   # Mill
GREEN    = RGBColor(0x05, 0x96, 0x69)   # Fabricator
AMBER    = RGBColor(0xd9, 0x77, 0x06)   # Integrator
GRAY     = RGBColor(0x6b, 0x72, 0x80)   # Auditor / 상태
WHITE    = RGBColor(0xff, 0xff, 0xff)
LIGHT    = RGBColor(0xe2, 0xe8, 0xf0)
PANEL    = RGBColor(0x12, 0x1a, 0x2e)   # 패널 배경
ACCENT   = RGBColor(0x38, 0xbd, 0xf8)   # 강조 (sky-400)
BORDER   = RGBColor(0x1e, 0x3a, 0x5f)   # 테두리

# ── 헬퍼 ──────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, border=None, radius=False):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=10, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_label_box(slide, label, l, t, w, h, bg, text_color=WHITE, size=9):
    r = add_rect(slide, l, t, w, h, fill=bg)
    add_text_box(slide, label, l + 0.03, t + 0.01, w - 0.06, h - 0.02,
                 font_size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return r


def arrow_right(slide, x, y, length=0.4, color=ACCENT):
    """가로 화살표"""
    from pptx.util import Pt
    from pptx.enum.shapes import PP_PLACEHOLDER
    # 선
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + length), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    # 화살촉 (작은 삼각형으로 표현)
    tip = slide.shapes.add_shape(5,  # right-triangle
        Inches(x + length - 0.05), Inches(y - 0.05),
        Inches(0.1), Inches(0.1))
    tip.fill.solid(); tip.fill.fore_color.rgb = color
    tip.line.fill.background()


def arrow_down(slide, x, y, length=0.25, color=ACCENT):
    """세로 화살표"""
    conn = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x), Inches(y + length))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)


# ── 슬라이드 생성 ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)

# ── 배경 ──────────────────────────────────────────────────────────────────────
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = BG

# ═══════════════════════════════════════════════════════════════════════════════
#  TITLE 영역
# ═══════════════════════════════════════════════════════════════════════════════
add_rect(slide, 0, 0, 13.33, 0.72, fill=PANEL, border=BORDER)
add_text_box(slide, "MTC Chain  —  강재 이력 관리 시스템 전체 플로우",
             0.2, 0.03, 9, 0.45, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "Ethereum Sepolia  ·  OpenZeppelin AccessControl  ·  IPFS/Pinata  ·  Cloudflare Workers",
             0.2, 0.42, 9, 0.28, font_size=9, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

# 오른쪽 상단 배지
add_label_box(slide, "Solidity 0.8.24", 10.5, 0.08, 1.35, 0.24, BLUE, size=8)
add_label_box(slide, "React 18 + ethers.js v6", 10.5, 0.38, 1.80, 0.24, RGBColor(0x5b, 0x21, 0xb6), size=8)

# ═══════════════════════════════════════════════════════════════════════════════
#  3-COLUMN LAYOUT
#  Col A: 역할 (0 ~ 2.5)
#  Col B: 비즈니스 플로우 (2.6 ~ 8.5)
#  Col C: 기술 스택 (8.6 ~ 13.2)
# ═══════════════════════════════════════════════════════════════════════════════

# ─── COL A: 역할 배지 ──────────────────────────────────────────────────────────
add_text_box(slide, "참여자 (Actors)", 0.12, 0.80, 2.3, 0.25,
             font_size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

roles = [
    ("Admin",        "역할 부여·회수\nDEFAULT_ADMIN_ROLE",  NAVY,  0x97),
    ("Mill  (제강사)",  "MTC 발행\nMILL_ROLE",               BLUE,  0x97),
    ("Fabricator\n(가공사)", "분할 / 조합\nFABRICATOR_ROLE", GREEN, 0x97),
    ("Integrator\n(통합사)", "사용 등록\nINTEGRATOR_ROLE",   AMBER, 0x97),
    ("Auditor\n(감사)",      "이력 조회 (누구나)",             GRAY,  0x97),
]

role_y = [1.08, 2.05, 3.02, 4.15, 5.20]
for i, (name, desc, color, _) in enumerate(roles):
    add_rect(slide, 0.12, role_y[i], 2.3, 0.82, fill=color, border=BORDER)
    add_text_box(slide, name, 0.15, role_y[i] + 0.04, 2.24, 0.3,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, 0.15, role_y[i] + 0.38, 2.24, 0.4,
                 font_size=7.5, color=LIGHT, align=PP_ALIGN.CENTER)

# ─── COL B: 비즈니스 플로우 ─────────────────────────────────────────────────────
add_text_box(slide, "비즈니스 플로우 (Business Flow)", 2.6, 0.80, 5.8, 0.25,
             font_size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

flow_boxes = [
    # (label, sublabel, y, bg_color, left_role_color)
    ("① 역할 부여",
     "grantMill() / grantFabricator() / grantIntegrator()\n→ AccessControl에 역할 등록, RoleGranted 이벤트 발생",
     1.08, NAVY),
    ("② MTC 발행  (issueMtc)",
     "PDF 해시(SHA-256) 계산 → Pinata에 PDF 업로드(IPFS CID 획득)\n→ issueMtc(steelId, weight, ipfsCid, pdfHash) 호출 → SteelMinted 이벤트",
     2.05, BLUE),
    ("③ 소유권 이전  (transferOwnership)",
     "transferOwnership(steelId, to) → ACTIVE 강재의 owner 변경\n→ SteelOwnershipTransferred 이벤트 발생",
     3.02, RGBColor(0x1d, 0x4e, 0xd8)),
    ("④-a 분할  splitSteel\n④-b 조합  combineSteel",
     "splitSteel(parentId, childWeights[]) : 부모 → SPLIT, 자식 ACTIVE 생성\ncombineSteel(parentIds[], childId, weight, cid, hash) : 부모들 → COMBINED",
     4.05, GREEN),
    ("⑤ 사용 등록  (markAsUsed)",
     "markAsUsed(steelId, productId) → 상태 ACTIVE → USED\nproductMap[productId] = steelId 저장, SteelUsed 이벤트",
     5.20, AMBER),
    ("⑥ 이력 조회  (Ancestry Tree)",
     "getSteel() / getParents() / getChildren() / getSteelByProduct()\n재귀 트리 탐색 → React Flow 시각화 + PDF 원본 IPFS 검증",
     6.20, GRAY),
]

for label, sub, fy, col in flow_boxes:
    add_rect(slide, 2.62, fy, 5.76, 0.82, fill=PANEL, border=col)
    # 왼쪽 컬러 바
    add_rect(slide, 2.62, fy, 0.08, 0.82, fill=col)
    add_text_box(slide, label, 2.75, fy + 0.04, 5.5, 0.28,
                 font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, sub, 2.75, fy + 0.35, 5.5, 0.44,
                 font_size=7.5, color=LIGHT, align=PP_ALIGN.LEFT)

# 플로우 화살표 (박스 사이)
arrow_ys = [1.92, 2.88, 3.88, 5.02, 6.03]
for ay in arrow_ys:
    arrow_down(slide, 5.5, ay, 0.14)

# 분기 표시 (④-a / ④-b)
add_text_box(slide, "OR", 5.35, 4.56, 0.3, 0.2,
             font_size=7, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ─── COL C: 기술 스택 ─────────────────────────────────────────────────────────
add_text_box(slide, "기술 스택 (Tech Stack)", 8.7, 0.80, 4.4, 0.25,
             font_size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

tech_layers = [
    # (title, detail, y, color)
    ("MetaMask Wallet",
     "사용자 서명(sign)\nEIP-712 트랜잭션 승인",
     1.08, RGBColor(0xe2, 0x52, 0x28)),

    ("Frontend\nReact 18 + Vite 5 + Tailwind 3",
     "ethers.js v6  Contract 호출\nReact Flow v11  Ancestry Tree\nCloudflare Pages  호스팅",
     1.90, RGBColor(0x06, 0x82, 0xf5)),

    ("Cloudflare Workers  (API Proxy)",
     "CORS 처리 + 비밀키 보호\nPinata API 호출 중계\nKV  productId 메타데이터 저장",
     3.10, RGBColor(0xf4, 0x81, 0x20)),

    ("Ethereum Sepolia  (Smart Contract)",
     "MtcRegistry.sol — AccessControl\n주요 함수: issueMtc · transferOwnership\n         splitSteel · combineSteel · markAsUsed\n이벤트: SteelMinted · SteelSplit · SteelCombined · SteelUsed",
     4.05, RGBColor(0x62, 0x7e, 0xea)),

    ("IPFS / Pinata  (분산 저장소)",
     "MTC PDF 원본 업로드 → CID 반환\nonchain pdfHash(SHA-256) 검증\n불변 이력 보장",
     5.55, RGBColor(0x11, 0x9e, 0x7e)),
]

for title, detail, ty, col in tech_layers:
    h = 0.75 if "Smart" in title else 0.65
    if "Smart" in title: h = 1.05
    add_rect(slide, 8.72, ty, 4.38, h, fill=PANEL, border=col)
    add_rect(slide, 8.72, ty, 0.07, h, fill=col)
    add_text_box(slide, title, 8.84, ty + 0.04, 4.2, 0.28,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, detail, 8.84, ty + 0.30, 4.2, h - 0.32,
                 font_size=7.5, color=LIGHT, align=PP_ALIGN.LEFT)

# 스택 화살표 (상위 → 하위 위임)
stack_arrow_ys = [1.74, 2.56, 3.76, 5.10, 6.21]
for say in stack_arrow_ys:
    arrow_down(slide, 10.91, say, 0.13, color=RGBColor(0x4a, 0x5e, 0x7a))

# ─── 하단 주석 바 ────────────────────────────────────────────────────────────
add_rect(slide, 0, 7.1, 13.33, 0.4, fill=PANEL, border=BORDER)
notes = (
    "상태 머신:  ACTIVE  →  SPLIT / COMBINED / USED  (불가역)  "
    "  |  무게 보존:  분할 손실 ≤ 10%,  조합 손실 ≤ 15%  "
    "  |  온체인 검증:  pdfHash(bytes32) = keccak256 or SHA-256  "
    "  |  역할 검증:  onlyRole() → AccessControlUnauthorizedAccount revert"
)
add_text_box(slide, notes, 0.15, 7.14, 13.0, 0.32,
             font_size=7.5, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
#  저장
# ═══════════════════════════════════════════════════════════════════════════════
out = r"C:\Users\wltn4\Desktop\MTC_Project\MTC_SystemFlow_Slide.pptx"
prs.save(out)
print(f"저장 완료: {out}")
